import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
from pathlib import Path
import zipfile
import re
import logging
import base64

logger = logging.getLogger(__name__)

# Configuraciones de límites
MAX_PAGINAS_PDF = 5
MAX_PARRAFOS_DOCX = 50
MAX_SLIDES_PPTX = 10
MAX_FILAS_XLSX = 50
MAX_TEXTO_CHARS = 2000

def limpiar_xml(xml_content):
    """Elimina etiquetas XML para dejar solo texto."""
    return re.sub(r'<[^>]+>', ' ', xml_content)

def leer_docx_raw(ruta):
    try:
        with zipfile.ZipFile(ruta) as z:
            xml_content = z.read('word/document.xml').decode('utf-8')
            return limpiar_xml(xml_content)
    except Exception:
        return ""

def leer_pptx_raw(ruta):
    texto = ""
    try:
        with zipfile.ZipFile(ruta) as z:
            slides = [f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
            for slide in slides:
                xml_content = z.read(slide).decode('utf-8')
                texto += limpiar_xml(xml_content) + " "
    except Exception:
        pass
    return texto

def leer_pdf(ruta):
    texto = ""
    with fitz.open(ruta) as doc:
        for pagina in doc.pages(0, min(doc.page_count, MAX_PAGINAS_PDF)): 
            texto += pagina.get_text()
    return texto

def leer_docx(ruta):
    try:
        doc = docx.Document(ruta)
        parrafos = [p.text for p in doc.paragraphs[:MAX_PARRAFOS_DOCX]]
        return "\n".join(parrafos)
    except Exception:
        return leer_docx_raw(ruta)

def leer_pptx(ruta):
    texto = ""
    try:
        prs = Presentation(ruta)
        for slide in prs.slides[:MAX_SLIDES_PPTX]:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + " "
    except Exception:
        texto = leer_pptx_raw(ruta)
    return texto

def leer_xlsx(ruta):
    texto = ""
    wb = openpyxl.load_workbook(ruta, data_only=True, read_only=True)
    sheet = wb.active
    for row in sheet.iter_rows(max_row=MAX_FILAS_XLSX, max_col=10, values_only=True):
        fila_texto = " ".join([str(cell) for cell in row if cell is not None])
        texto += fila_texto + " "
    return texto

def leer_texto_plano(ruta):
    try:
        with open(ruta, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read(MAX_TEXTO_CHARS)
    except Exception:
        return ""

def encode_image(image_path):
    """Codifica una imagen en base64."""
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

# Registro de lectores por extensión
LECTORES = {
    '.pdf': leer_pdf,
    '.docx': leer_docx,
    '.pptx': leer_pptx,
    '.xlsx': leer_xlsx,
    '.txt': leer_texto_plano,
    '.csv': leer_texto_plano,
    '.md': leer_texto_plano,
    '.json': leer_texto_plano,
    '.xml': leer_texto_plano,
    '.html': leer_texto_plano,
}

# Extensiones de imagen soportadas
EXTENSIONES_IMAGEN = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff'}

def extraer_texto_del_archivo(ruta: Path):
    """
    Función principal para extraer texto de un archivo según su extensión.
    """
    extension = ruta.suffix.lower()
    
    if extension in LECTORES:
        try:
            return LECTORES[extension](ruta).lower().strip()
        except Exception as e:
            logger.error(f"Error al leer documento {ruta.name}: {e}")
            return ""
    
    if extension in EXTENSIONES_IMAGEN:
        return "IMAGE_FILE" # Marcador para que el servicio sepa que debe usar visión
        
    return ""

def es_extension_soportada(ruta: Path):
    ext = ruta.suffix.lower()
    return ext in LECTORES or ext in EXTENSIONES_IMAGEN