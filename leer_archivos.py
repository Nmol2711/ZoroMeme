import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
from pathlib import Path

def extraer_texto_del_archivo(ruta : Path):
    extension = ruta.suffix.lower()
    texto_extraido = ""

    try:
        # --- Leer PDF ---
        if extension == '.pdf':
            with fitz.open(ruta) as doc:
                # Leemos solo las primeras páginas para eficiencia
                for pagina in doc.pages(0, 3): 
                    texto_extraido += pagina.get_text()

        # --- Leer WORD (.docx) ---
        elif extension == '.docx':
            doc = docx.Document(ruta)
            parrafos = [p.text for p in doc.paragraphs[:30]]
            texto_extraido = "\n".join(parrafos)

        # --- Leer POWERPOINT (.pptx) ---
        elif extension == '.pptx':
            prs = Presentation(ruta)
            # Leemos solo las primeras 10 diapositivas para no saturar
            for i, slide in enumerate(prs.slides[:10]):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto_extraido += shape.text + " "

        # --- Leer EXCEL (.xlsx) ---
        elif extension == '.xlsx':
            # data_only=True para leer el resultado de fórmulas, no la fórmula en sí
            wb = openpyxl.load_workbook(ruta, data_only=True, read_only=True)
            sheet = wb.active # Leemos la hoja principal
            # Leemos las primeras 50 filas y 10 columnas
            for row in sheet.iter_rows(max_row=50, max_col=10, values_only=True):
                fila_texto = " ".join([str(cell) for cell in row if cell is not None])
                texto_extraido += fila_texto + " "

    except Exception as e:
        print(f"Error al leer {ruta.name}: {e}")

    return texto_extraido.lower().strip()